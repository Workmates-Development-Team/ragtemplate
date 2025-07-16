import os
import json
import logging
import tempfile
import numpy as np
import boto3
import psycopg2
import whisper
import docx
import PyPDF2
import pandas as pd
import pytesseract
import ssl
import certifi
import uuid
from pptx import Presentation
from PIL import Image
from flask import Flask, request, jsonify, Response, send_from_directory
from flask_cors import CORS
from moviepy import AudioFileClip
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Logger setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# AWS & DB Config
BEDROCK_REGION = os.getenv("AWS_REGION", "us-east-1")
AWS_ACCESS_KEY = os.getenv("AWS_ACCESS_KEY")
AWS_SECRET_ACCESS_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
MODEL_ID_EMBED = os.getenv("MODEL_ID_EMBED")
MODEL_ID_CHAT = os.getenv("MODEL_ID_CHAT")
DATABASE_URL = os.getenv("DATABASE_URL")

# Whisper model
logger.info("Loading Whisper model...")
whisper_model = whisper.load_model("base")
logger.info("Whisper model loaded successfully.")

# AWS Bedrock client
bedrock = boto3.client(
    "bedrock-runtime",
    region_name=BEDROCK_REGION,
    aws_access_key_id=AWS_ACCESS_KEY,
    aws_secret_access_key=AWS_SECRET_ACCESS_KEY
)
logger.info("Bedrock client initialized.")

ssl._create_default_https_context = ssl._create_unverified_context
# If you want to use certifi's CA bundle instead, comment the above and use:
# ssl._create_default_https_context = lambda: ssl.create_default_context(cafile=certifi.where())

# Flask app setup
app = Flask(__name__)
CORS(app)

# --- Utilities ---
def get_db_connection():
    conn = psycopg2.connect(DATABASE_URL)
    logger.info("Connected to the database")
    return conn, conn.cursor()

def split_text(text, max_words=500):
    words = text.split()
    for i in range(0, len(words), max_words):
        yield ' '.join(words[i:i+max_words])

def transcribe_audio(file_path):
    result = whisper_model.transcribe(file_path)
    return result['text']

import base64

def embed_text_or_image(content, content_type='text', model_id=None):
    """
    Embed text or base64 image using Bedrock
    content_type: 'text' or 'image'
    """
    model_id = model_id or MODEL_ID_EMBED

    # Hard limit on content length to avoid token overflow
    max_chars = 16000  # ~4000 tokens safe margin
    if content_type == 'text' and len(content) > max_chars:
        logger.warning(f"Truncating content from {len(content)} chars to {max_chars} chars before embedding.")
        content = content[:max_chars]

    # Log chunk size before embedding
    logger.info(f"Embedding chunk of length {len(content)} characters")

    if content_type == 'text':
        request_body = json.dumps({
            "inputText": content,
            "dimensions": 256,
            "normalize": True
        })
    elif content_type == 'image':
        request_body = json.dumps({
            "inputImage": content,
            "dimensions": 256,
            "normalize": True
        })
    else:
        raise ValueError("Invalid content_type for embedding")

    response = bedrock.invoke_model(
        modelId=model_id,
        body=request_body,
        accept="application/json",
        contentType="application/json"
    )
    response_body = json.loads(response.get("body").read())
    return response_body.get("embedding")

def retrieve_similar_chunks(query_embedding, top_k=3):
    conn, cursor = get_db_connection()
    cursor.execute("SELECT chunk_index, embedding, file_name, file_path, chunk_text FROM embeddings")
    rows = cursor.fetchall()
    cursor.close()
    conn.close()

    similarities = []
    query_vec = np.array(query_embedding)

    for chunk_index, embedding, file_name, file_path, chunk_text in rows:
        if embedding is None:
            continue
        emb_vec = np.array(embedding)
        sim = np.dot(emb_vec, query_vec) / (np.linalg.norm(emb_vec) * np.linalg.norm(query_vec))
        similarities.append((sim, chunk_index, file_name, file_path, chunk_text))

    similarities.sort(key=lambda x: x[0], reverse=True)
    return similarities[:top_k]

def extract_text(file):
    filename = file.filename
    save_path = os.path.join(UPLOAD_DIR, filename)
    file.save(save_path)
    logger.info(f"Saved file at: {save_path}")

    suffix = os.path.splitext(filename)[1].lower()

    try:
        if suffix in ['.mp3', '.wav', '.m4a']:
            return {
                "text": transcribe_audio(save_path),
                "image_base64": None,
                "filename": filename,
                "file_path": f"/uploads/{filename}"
            }

        elif suffix in ['.mp4', '.mov', '.mkv']:
            audio_path = save_path + "_audio.wav"
            clip = AudioFileClip(save_path)
            clip.write_audiofile(audio_path)
            clip.close()
            return {
                "text": transcribe_audio(audio_path),
                "image_base64": None,
                "filename": filename,
                "file_path": f"/uploads/{filename}"
            }

        elif suffix in ['.jpg', '.jpeg', '.png']:
            try:
                with open(save_path, 'rb') as img_file:
                    image_bytes = bytearray(img_file.read())

                textract = boto3.client(
                    'textract',
                    region_name=BEDROCK_REGION,
                    aws_access_key_id=AWS_ACCESS_KEY,
                    aws_secret_access_key=AWS_SECRET_ACCESS_KEY
                )

                response = textract.detect_document_text(Document={'Bytes': image_bytes})

                extracted_lines = [
                    item["Text"] for item in response["Blocks"]
                    if item["BlockType"] == "LINE"
                ]
                textract_text = "\n".join(extracted_lines)

                return {
                    "text": textract_text,
                    "image_base64": None,
                    "filename": filename,
                    "file_path": f"/uploads/{filename}"
                }
            except Exception as e:
                if "UnrecognizedClientException" in str(e) or "security token" in str(e):
                    logger.error("AWS credentials are invalid or missing for Textract.")
                logger.error(f"Textract OCR failed: {e}")
                raise

        elif suffix.endswith('.pdf'):
            reader = PyPDF2.PdfReader(save_path)
            text = '\n'.join([page.extract_text() for page in reader.pages if page.extract_text()])
            return {
                "text": text,
                "image_base64": None,
                "filename": filename,
                "file_path": f"/uploads/{filename}"
            }

        elif suffix.endswith('.docx'):
            doc = docx.Document(save_path)
            full_text = []
            full_text.extend([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    full_text.append(row_text)
            return {
                "text": '\n'.join(full_text),
                "image_base64": None,
                "filename": filename,
                "file_path": f"/uploads/{filename}"
            }

        elif suffix.endswith('.pptx'):
            try:
                presentation = Presentation(save_path)
                full_text = []
                for slide in presentation.slides:
                    # Extract text from shapes (e.g., text boxes, placeholders)
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            full_text.append(shape.text.strip())
                    # Extract text from tables
                    for shape in slide.shapes:
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                                full_text.append(row_text)
                return {
                    "text": '\n'.join(full_text),
                    "image_base64": None,
                    "filename": filename,
                    "file_path": f"/uploads/{filename}"
                }
            except Exception as e:
                logger.error(f"Error reading .pptx file: {e}")
                raise

        elif suffix.endswith('.txt') or suffix.endswith('.py'):
            with open(save_path, 'r', encoding='utf-8') as f:
                return {
                    "text": f.read(),
                    "image_base64": None,
                    "filename": filename,
                    "file_path": f"/uploads/{filename}"
                }

        elif suffix.endswith('.csv'):
            df = pd.read_csv(save_path)
            return {
                "text": df.to_string(index=False),
                "image_base64": None,
                "filename": filename,
                "file_path": f"/Uploads/{filename}"
            }

        elif suffix.endswith('.xlsx'):
            try:
                df = pd.read_excel(save_path, engine="openpyxl")
                return {
                    "text": df.to_string(index=False),
                    "image_base64": None,
                    "filename": filename,
                    "file_path": f"/Uploads/{filename}"
                }
            except Exception as e:
                logger.error(f"Error reading .xlsx file: {e}")
                raise

        else:
            raise ValueError("Unsupported file type")

    finally:
        try:
            # os.remove(save_path)
            logger.info(f"File retained at: {save_path}")
        except Exception as e:
            logger.warning(f"Failed to delete temp file: {e}")

# --- Routes ---

@app.route('/upload', methods=['POST'])
def upload_file():
    logger.info("File upload endpoint called")
    if 'file' not in request.files:
        return jsonify({'error': 'No file(s) uploaded'}), 400

    files = request.files.getlist('file')
    results = []

    conn, cursor = get_db_connection()
    logger.info("Connected to the database for embedding")
    cursor.execute('''CREATE TABLE IF NOT EXISTS embeddings (
        id SERIAL PRIMARY KEY,
        file_name TEXT,
        file_path TEXT,
        chunk_index INTEGER,
        embedding FLOAT8[],
        embedding_size INTEGER,
        chunk_text TEXT
    )''')
    conn.commit()

    for file in files:
        upload_id = uuid.uuid4().hex
        # try:
        #     extracted = extract_text(file)
        #     text = extracted.get("text", "")
        #     image_base64 = extracted.get("image_base64", None)
        #     filename = extracted["filename"]
        #     file_path = extracted["file_path"]
        #     logger.info(f"Extracted content from {filename}")
        # except Exception as e:
        #     logger.error(f"Text/Image extraction failed for {file.filename}: {e}")
        #     results.append({'filename': file.filename, 'error': 'Failed to extract content'})
        #     continue

        # chunk_count = 0
        # if text:
        #     logger.info(f"Embedding text chunks for {filename}")
        #     for idx, chunk in enumerate(split_text(text)):
        #         embedding = embed_text_or_image(chunk, content_type='text')
        #         cursor.execute(
        #             "INSERT INTO embeddings (file_name, file_path, chunk_index, embedding, embedding_size, chunk_text) VALUES (%s, %s, %s, %s, %s, %s)",
        #             (filename, file_path, idx + 1, embedding, len(embedding), chunk)
        #         )
        #         chunk_count += 1
        try:
            extracted = extract_text(file)
            text = extracted.get("text", "")
            image_base64 = extracted.get("image_base64", None)
            filename = extracted["filename"]
            file_path = extracted["file_path"]
            logger.info(f"Extracted content from {filename} with upload_id {upload_id}")
        except Exception as e:
            logger.error(f"Text/Image extraction failed for {file.filename}: {e}")
            results.append({'filename': file.filename, 'error': 'Failed to extract content'})
            continue

        chunk_count = 0
        if text:
            for idx, chunk in enumerate(split_text(text)):
                embedding = embed_text_or_image(chunk, content_type='text')
                cursor.execute(
                    "INSERT INTO embeddings (file_name, file_path, chunk_index, embedding, embedding_size, chunk_text, upload_id) VALUES (%s, %s, %s, %s, %s, %s, %s)",
                    (filename, file_path, idx + 1, embedding, len(embedding), chunk, upload_id)
                )
                chunk_count += 1

        if image_base64:
            try:
                logger.info(f"Embedding image for {filename}")
                embedding = embed_text_or_image(image_base64, content_type='image', model_id="amazon.nova-pro-v1:0")
                cursor.execute(
                    "INSERT INTO embeddings (file_name, file_path, chunk_index, embedding, embedding_size, chunk_text) VALUES (%s, %s, %s, %s, %s, %s)",
                    (filename, file_path, 9999, embedding, len(embedding), "[IMAGE_CONTENT]")
                )
                chunk_count += 1
            except Exception as e:
                logger.warning(f"Image embedding failed: {e}")

        conn.commit()
        results.append({'filename': filename, 'status': 'success', 'chunks_processed': chunk_count})

    cursor.close()
    conn.close()

    logger.info(f"Processed files: {', '.join([r['filename'] for r in results])}")
    return jsonify({'message': 'Processed files', 'results': results})
@app.route('/list_files', methods=['GET'])
def list_files():
    try:
        conn, cursor = get_db_connection()
        cursor.execute("""
            SELECT file_name, file_path, upload_id, COUNT(*) as chunk_count
            FROM embeddings
            GROUP BY file_name, file_path, upload_id
        """)
        files = [
            {
                "file_name": row[0],
                "file_path": row[1],
                "upload_id": row[2],
                "chunk_count": row[3]
            }
            for row in cursor.fetchall()
        ]
        cursor.close()
        conn.close()
        logger.info(f"Retrieved {len(files)} unique file entries")
        return jsonify({'files': files})
    except Exception as e:
        logger.error(f"Failed to list files: {e}")
        return jsonify({'error': 'Failed to retrieve files'}), 500
@app.route('/chat', methods=['POST'])
def chat():
    data = request.get_json()
    user_query = data.get('message') or data.get('question')
    model_id = data.get('model_id', MODEL_ID_CHAT)

    if not user_query:
        return jsonify({'error': 'Missing question'}), 400

    try:
        query_embedding = embed_text_or_image(user_query, content_type='text')
        top_chunks = retrieve_similar_chunks(query_embedding)

        context = "\n\n".join([chunk_text for _, _, _, _, chunk_text in top_chunks])
        references = [(file_name, file_path, chunk_index) for _, chunk_index, file_name, file_path, _ in top_chunks]
        
        prompt = f"Context:\n{context}\n\nQuestion: {user_query}\nAnswer:"

        request_body = {
            "schemaVersion": "messages-v1",
            "messages": [{"role": "user", "content": [{"text": prompt}]}],
            "system": [{"text": "You are a helpful assistant. Use the provided context to answer the user's question as accurately as possible. Respond by stating the answer clearly, no additional explanation. Do not use any external data beyond the provided context."}],
            "inferenceConfig": {
                "maxTokens": 512,
                "topP": 0.9,
                "topK": 20,
                "temperature": 0.7
            }
        }

        response = bedrock.invoke_model_with_response_stream(
            modelId=model_id,
            body=json.dumps(request_body)
        )

        def stream_response():
            stream = response.get("body")
            answer = ""
            if stream:
                for event in stream:
                    chunk = event.get("chunk")
                    if chunk:
                        chunk_json = json.loads(chunk.get("bytes").decode())
                        delta = chunk_json.get("contentBlockDelta", {}).get("delta", {}).get("text", "")
                        answer += delta

            references_json = [
                {
                    "file_name": file_name,
                    "file_path": file_path,
                    "chunk_index": chunk_index
                }
                for idx, (file_name, file_path, chunk_index) in enumerate(references, 1)
            ]
            # Yield the whole response as a single JSON object
            yield json.dumps({
                "answer": answer,
                "references": references_json
            }, ensure_ascii=False)

        return Response(stream_response(), content_type='application/json')

    except Exception as e:
        logger.error(f"Chat failed: {e}")
        return jsonify({'error': 'Chat failed'}), 500

@app.route('/delete', methods=['POST'])
def delete_upload():
    data = request.get_json()
    upload_id = data.get('upload_id')

    if not upload_id:
        return jsonify({'error': 'upload_id is required'}), 400

    try:
        conn, cursor = get_db_connection()

        # Get associated file paths
        cursor.execute("SELECT DISTINCT file_path FROM embeddings WHERE upload_id = %s", (upload_id,))
        file_paths = cursor.fetchall()

        # Delete rows from database
        cursor.execute("DELETE FROM embeddings WHERE upload_id = %s", (upload_id,))
        conn.commit()

        # Delete files from filesystem
        for (file_path,) in file_paths:
            local_path = os.path.join(".", file_path.lstrip("/"))
            try:
                if os.path.exists(local_path):
                    os.remove(local_path)
                    logger.info(f"Deleted file: {local_path}")
                else:
                    logger.warning(f"File not found: {local_path}")
            except Exception as e:
                logger.warning(f"Failed to delete file {local_path}: {e}")

        cursor.close()
        conn.close()

        logger.info(f"Deleted upload with ID: {upload_id}")
        return jsonify({'message': f'Deletion completed for upload_id: {upload_id}'})

    except Exception as e:
        logger.error(f"Deletion failed for upload_id {upload_id}: {e}")
        return jsonify({'error': 'Deletion failed'}), 500


# Directory for uploaded files
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

@app.route('/uploads/<path:filename>')
def serve_upload(filename):
    return send_from_directory(UPLOAD_DIR, filename)

# Run the app
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)